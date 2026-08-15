#!/usr/bin/env python3
"""Extract content and images from PowerPoint file"""

import argparse
import hashlib
import os
import sys

# Try to import pptx, install if needed
try:
    from pptx import Presentation
except ImportError:
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "Pillow", "-q"])
    from pptx import Presentation

import json
import posixpath
import re
import zipfile
from pathlib import Path
from urllib.request import urlopen
from xml.etree import ElementTree as ET

from pptx.enum.shapes import MSO_SHAPE_TYPE


def iter_shapes(shape):
    yield shape
    children = getattr(shape, "shapes", None)
    if children:
        for child in children:
            yield from iter_shapes(child)


def extract_images_from_shape(shape, slide_number, output_dir, slide_content, image_prefix):
    extracted = False
    for candidate in iter_shapes(shape):
        if getattr(candidate, "shape_type", None) != MSO_SHAPE_TYPE.PICTURE:
            image = getattr(candidate, "image", None)
            if image is None:
                continue
        else:
            image = getattr(candidate, "image", None)
            if image is None:
                continue

        try:
            image_hash = hashlib.sha256(image.blob).hexdigest()
            if image_hash in slide_content.setdefault("_image_hashes", set()):
                continue
            slide_content["_image_hashes"].add(image_hash)
            image_filename = f"{image_prefix}_slide{slide_number}_image{len(slide_content['images'])}.png"
            image_path = output_dir / image_filename

            with open(image_path, 'wb') as f:
                f.write(image.blob)

            slide_content["images"].append({
                "filename": image_filename,
                "path": f"images/{image_filename}"
            })
            print(f"  Saved image: {image_filename}")
            extracted = True
        except Exception as e:
            print(f"  Error saving image: {e}")

    return extracted


def extract_images_from_package(slide_number, pptx_path, output_dir, slide_content, image_prefix):
    extracted = False
    try:
        with zipfile.ZipFile(pptx_path) as archive:
            slide_xml_name = f"ppt/slides/slide{slide_number}.xml"
            if slide_xml_name not in archive.namelist():
                return False

            slide_xml = archive.read(slide_xml_name).decode("utf-8", errors="ignore")
            rels_path = f"ppt/slides/_rels/slide{slide_number}.xml.rels"
            rels_xml = archive.read(rels_path).decode("utf-8", errors="ignore") if rels_path in archive.namelist() else ""
            if not rels_xml:
                return False

            rels = ET.fromstring(rels_xml.encode("utf-8"))
            rel_map = {}
            for rel in rels.findall('{http://schemas.openxmlformats.org/package/2006/relationships}Relationship'):
                rel_id = rel.attrib.get('Id')
                if rel_id:
                    rel_map[rel_id] = {
                        'target': rel.attrib.get('Target', ''),
                        'target_mode': rel.attrib.get('TargetMode', '')
                    }

            seen_targets = set()
            for match in re.finditer(r'<a:blip[^>]*(?:r:embed|r:link)="([^"]+)"', slide_xml):
                rel_id = match.group(1)
                rel_info = rel_map.get(rel_id)
                if not rel_info:
                    continue

                target = rel_info['target']
                target_mode = rel_info['target_mode']
                lookup_key = f"{rel_id}:{target}"
                if lookup_key in seen_targets:
                    continue
                seen_targets.add(lookup_key)

                try:
                    if target_mode == 'External' or target.startswith(('http://', 'https://')):
                        image_bytes = urlopen(target).read()
                        ext = '.bin'
                    else:
                        source_dir = posixpath.dirname(slide_xml_name)
                        resolved_path = posixpath.normpath(posixpath.join(source_dir, target))
                        if resolved_path.startswith('/'):
                            resolved_path = resolved_path.lstrip('/')
                        image_bytes = archive.read(resolved_path)
                        ext = posixpath.splitext(resolved_path)[1] or '.bin'

                    image_hash = hashlib.sha256(image_bytes).hexdigest()
                    if image_hash in slide_content.setdefault("_image_hashes", set()):
                        continue
                    slide_content["_image_hashes"].add(image_hash)
                    image_filename = f"{image_prefix}_slide{slide_number}_image{len(slide_content['images'])}{ext or '.bin'}"
                    image_path = output_dir / image_filename
                    with open(image_path, 'wb') as f:
                        f.write(image_bytes)

                    slide_content["images"].append({
                        "filename": image_filename,
                        "path": f"images/{image_filename}"
                    })
                    print(f"  Saved image: {image_filename}")
                    extracted = True
                except Exception as e:
                    print(f"  Error saving image: {e}")
    except Exception as e:
        print(f"  Error reading slide package: {e}")

    return extracted


def build_args():
    parser = argparse.ArgumentParser(description="Extract content and images from a lecture PowerPoint file.")
    parser.add_argument(
        "--lecture-id",
        default="lect-01a",
        help="Lecture id used for output naming, e.g. lect-01b",
    )
    parser.add_argument(
        "--pptx",
        default="3DIMD Lecture 01a Module Info 01.pptx",
        help="PowerPoint filename inside the source PPTX directory.",
    )
    parser.add_argument(
        "--source-dir",
        default="site/source-lecture-pptx",
        help="Directory containing lecture PPTX source files.",
    )
    parser.add_argument(
        "--pages-dir",
        default="site/pages",
        help="Directory where lecture page assets and JSON are generated.",
    )
    parser.add_argument(
        "--images-subdir",
        default="images",
        help="Images subdirectory under pages directory.",
    )
    return parser.parse_args()


def main():
    args = build_args()
    repo_root = Path(__file__).parent
    pptx_path = repo_root / Path(args.source_dir) / args.pptx
    pages_dir = repo_root / Path(args.pages_dir)
    output_dir = pages_dir / args.images_subdir
    output_file = pages_dir / f"{args.lecture_id}-content.json"
    image_prefix = args.lecture_id.replace("/", "-").replace("\\", "-")

    output_dir.mkdir(parents=True, exist_ok=True)

    print(f"Extracting from: {pptx_path}")

    prs = Presentation(str(pptx_path))
    slides_data = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_content = {
            "slide_number": slide_num,
            "title": "",
            "section_title": "",
            "content": [],
            "images": []
        }

        # Package relationships belong to the slide, not each individual shape.
        extract_images_from_package(
            slide_num,
            pptx_path,
            output_dir,
            slide_content,
            image_prefix,
        )

        text_shapes = [
            shape for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]
        title_shape = next(
            (shape for shape in text_shapes if shape.name.startswith("Title")),
            None,
        )
        if title_shape is None and text_shapes:
            first_shape = text_shapes[0]
            first_text = first_shape.text.strip()
            if len(first_text) <= 120 and "\n" not in first_text:
                title_shape = first_shape
        title_shape_name = title_shape.name if title_shape is not None else None

        # Extract text and shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                if shape.name == title_shape_name:
                    slide_content["title"] = text
                else:
                    slide_content["content"].append(text)

            # Extract images embedded directly in shapes and nested groups.
            extract_images_from_shape(
                shape,
                slide_num,
                output_dir,
                slide_content,
                image_prefix,
            )

        if slide_content["content"]:
            first_block = slide_content["content"][0]
            block_lines = first_block.splitlines()
            first_line_index = next((index for index, line in enumerate(block_lines) if line.strip()), None)
            first_line = block_lines[first_line_index].strip() if first_line_index is not None else ""
            title_match = re.match(r"^(.*?)\s+[–-]\s+", first_line)
            section_title = title_match.group(1).strip() if title_match else first_line
            slide_content["section_title"] = section_title
            if first_line_index is not None:
                remove_title_line = not title_match
                remaining_lines = block_lines[first_line_index + 1:] if remove_title_line else block_lines[first_line_index:]
                if remaining_lines:
                    slide_content["content"][0] = "\n".join(remaining_lines)
                else:
                    slide_content["content"].pop(0)

        slide_content.pop("_image_hashes", None)
        slides_data.append(slide_content)
        print(f"Slide {slide_num}: {slide_content['title']}")

    with open(output_file, 'w', encoding='utf-8') as f:
        json.dump(slides_data, f, indent=2, ensure_ascii=False)

    print(f"\nExtracted {len(slides_data)} slides")
    print(f"Content saved to: {output_file}")
    print("\nSlide Summary:")
    for slide in slides_data:
        print(f"\nSlide {slide['slide_number']}: {slide['title']}")
        for line in slide['content']:
            print(f"  - {line[:80]}")
        if slide['images']:
            print(f"  Images: {len(slide['images'])}")


if __name__ == "__main__":
    main()
